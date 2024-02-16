import os
from pptx import Presentation
from pptx.util import Cm

# Ввод параметров
prs_name = input("Введите название презентации: ")
photo_folder = input("Введите название папки с фото: ")
slide_index = int(input("Введите индекс стартового слайда (нумерация начинается с 0): "))
save_path = input("Введите название сохраненной презентации: ")

# Открываем существующую презентацию
prs = Presentation(prs_name)

# Размеры фотографии
photo_width = Cm(15)  # Ширина фотографии
photo_height = Cm(8.2)  # Высота фотографии

# Координаты фотографии по вертикали и горизонтали
photo_coordinates = [
    (Cm(2.4), Cm(4.43)),  # Верхнее левое фото
    (Cm(18.09), Cm(4.43)),  # Верхнее правое фото
    (Cm(2.4), Cm(15.52)),  # Нижнее левое фото
    (Cm(18.09), Cm(15.52))  # Нижнее правое фото
]

# Ключевые слова для поиска
photo_format = ('jpeg', 'jpg', 'png')

# Получаем список всех фото из папки
photos = [photo for photo in os.listdir(photo_folder) if photo.endswith(photo_format)]

# Вычисляем количество слайдов, которые понадобятся для всех фото
num_slides = (len(photos) + 3) // 4

# Вставляем фотографии на слайды, начиная с заданного индекса слайда
photo_index = 0

for _ in range(num_slides):
    # Получаем очередной слайд
    slide = prs.slides[slide_index]

    # Вставляем фотографии на слайд
    for i, coord in enumerate(photo_coordinates):
        if photo_index >= len(photos):
            break

        left, top = coord
        photo_path = os.path.join(photo_folder, photos[photo_index])
        pic = slide.shapes.add_picture(photo_path, left, top, photo_width, photo_height)

        photo_index += 1

    slide_index += 1

# Сохраняем изменения
prs.save(save_path)
